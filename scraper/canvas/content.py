"""Turning Canvas HTML and attached files into text an LLM can actually use.

Two jobs:
  * HTML -> markdown, preserving tables (marking criteria are nearly always
    tables, and plain-text conversion flattens them into unusable runs)
  * attached files -> text, so an assignment brief that lives in a PDF is
    still readable from the JSON
"""

from __future__ import annotations

import re
from pathlib import Path
from urllib.parse import urljoin, urlparse

from bs4 import BeautifulSoup

MAX_EXTRACT_CHARS = 200_000

# Canvas stamps internal rich-content links with these attributes, which tells
# us what a link points at without having to guess from the URL shape.
API_TYPE_ATTR = "data-api-returntype"
API_ENDPOINT_ATTR = "data-api-endpoint"

_URL_KIND_PATTERNS = [
    (re.compile(r"/files/\d+"), "file"),
    (re.compile(r"/assignments/\d+"), "assignment"),
    (re.compile(r"/pages/"), "page"),
    (re.compile(r"/quizzes/\d+"), "quiz"),
    (re.compile(r"/discussion_topics/\d+"), "discussion"),
    (re.compile(r"/modules"), "module"),
]


def _absolutise(soup: BeautifulSoup, base_url: str | None) -> None:
    if not base_url:
        return
    for tag, attr in (("a", "href"), ("img", "src"), ("iframe", "src")):
        for node in soup.find_all(tag):
            value = node.get(attr)
            if value and not value.startswith(("http://", "https://", "mailto:", "#")):
                node[attr] = urljoin(base_url, value)


def html_to_markdown(html: str | None, base_url: str | None = None) -> str:
    """Convert Canvas rich content to markdown. Returns '' for empty input."""
    if not html or not html.strip():
        return ""

    soup = BeautifulSoup(html, "html.parser")
    for node in soup(["script", "style"]):
        node.decompose()
    _absolutise(soup, base_url)

    try:
        from markdownify import markdownify

        text = markdownify(str(soup), heading_style="ATX", bullets="-")
    except ImportError:
        text = soup.get_text("\n")

    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _classify(url: str, api_type: str | None) -> str:
    if api_type:
        return api_type.lower()
    for pattern, kind in _URL_KIND_PATTERNS:
        if pattern.search(url):
            return kind
    return "external"


def extract_links(html: str | None, base_url: str | None = None) -> list[dict]:
    """Every outbound link in a blob of Canvas HTML, deduped and classified."""
    if not html or not html.strip():
        return []

    soup = BeautifulSoup(html, "html.parser")
    _absolutise(soup, base_url)

    found: dict[str, dict] = {}

    for node in soup.find_all("a", href=True):
        url = node["href"].strip()
        if not url or url.startswith("#"):
            continue
        label = " ".join(node.get_text(" ", strip=True).split())[:200]
        entry = {
            "url": url,
            "label": label or None,
            "kind": _classify(url, node.get(API_TYPE_ATTR)),
        }
        if node.get(API_ENDPOINT_ATTR):
            entry["api_endpoint"] = node[API_ENDPOINT_ATTR]
        # Prefer the version that carries a human label.
        if url not in found or (entry["label"] and not found[url].get("label")):
            found[url] = entry

    for node in soup.find_all(["iframe", "embed"], src=True):
        url = node["src"].strip()
        if url and url not in found:
            found[url] = {
                "url": url,
                "label": node.get("title") or None,
                "kind": "embed",
                "host": urlparse(url).netloc or None,
            }

    return list(found.values())


# --------------------------------------------------------------- file -> text

def _pdf_text(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for number, page in enumerate(reader.pages, 1):
        body = (page.extract_text() or "").strip()
        if body:
            pages.append(f"[page {number}]\n{body}")
    return "\n\n".join(pages)


def _pptx_text(path: Path) -> str:
    from pptx import Presentation

    deck = Presentation(str(path))
    slides = []
    for number, slide in enumerate(deck.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
        # Speaker notes are often where the real explanation lives.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Notes: {notes}")
        if parts:
            slides.append(f"[slide {number}]\n" + "\n".join(parts))
    return "\n\n".join(slides)


def _docx_text(path: Path) -> str:
    import docx

    document = docx.Document(str(path))
    blocks = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                blocks.append(" | ".join(cells))
    return "\n".join(blocks)


def _xlsx_text(path: Path) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheets = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[sheet {sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def _plain_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


_EXTRACTORS = {
    ".pdf": ("pypdf", _pdf_text),
    ".pptx": ("python-pptx", _pptx_text),
    ".docx": ("python-docx", _docx_text),
    ".xlsx": ("openpyxl", _xlsx_text),
    ".txt": ("builtin", _plain_text),
    ".md": ("builtin", _plain_text),
    ".csv": ("builtin", _plain_text),
    ".json": ("builtin", _plain_text),
}



# --------------------------------------------------------------------- OCR

OCR_DPI = 200
OCR_MAX_PAGES = 40


def _ocr_pdf(path: Path, dpi: int = OCR_DPI, max_pages: int = OCR_MAX_PAGES) -> str:
    """OCR a scanned PDF using the OS text recogniser.

    Scanned readings have no text layer at all, so pypdf returns nothing and the
    file is invisible to anything reading the JSON. macOS ships a good enough
    recogniser in Vision, which keeps this to a pip install rather than a system
    one. Raises ImportError off macOS, which file_to_text turns into a note.
    """
    import Quartz
    import Vision
    from Foundation import NSData

    data = NSData.dataWithContentsOfFile_(str(path))
    provider = Quartz.CGDataProviderCreateWithCFData(data)
    document = Quartz.CGPDFDocumentCreateWithProvider(provider)
    if not document:
        raise ValueError("could not open PDF for rendering")

    total = Quartz.CGPDFDocumentGetNumberOfPages(document)
    scale = dpi / 72.0
    pages = []

    for number in range(1, min(total, max_pages) + 1):
        page = Quartz.CGPDFDocumentGetPage(document, number)
        rect = Quartz.CGPDFPageGetBoxRect(page, Quartz.kCGPDFMediaBox)
        width = int(rect.size.width * scale) or 1
        height = int(rect.size.height * scale) or 1
        context = Quartz.CGBitmapContextCreate(
            None, width, height, 8, 0,
            Quartz.CGColorSpaceCreateDeviceRGB(), Quartz.kCGImageAlphaNoneSkipLast,
        )
        # Scans are dark-on-white; without this the transparent ground reads black.
        Quartz.CGContextSetRGBFillColor(context, 1, 1, 1, 1)
        Quartz.CGContextFillRect(context, Quartz.CGRectMake(0, 0, width, height))
        Quartz.CGContextScaleCTM(context, scale, scale)
        Quartz.CGContextDrawPDFPage(context, page)
        image = Quartz.CGBitmapContextCreateImage(context)

        lines: list[str] = []

        def collect(request, error, _lines=lines):
            for observation in request.results() or []:
                best = observation.topCandidates_(1)
                if best:
                    _lines.append(best[0].string())

        request = Vision.VNRecognizeTextRequest.alloc().initWithCompletionHandler_(collect)
        request.setRecognitionLevel_(Vision.VNRequestTextRecognitionLevelAccurate)
        request.setUsesLanguageCorrection_(True)
        handler = Vision.VNImageRequestHandler.alloc().initWithCGImage_options_(image, None)
        handler.performRequests_error_([request], None)

        body = "\n".join(lines).strip()
        if body:
            pages.append(f"[page {number}]\n{body}")

    if total > max_pages:
        pages.append(f"[{total - max_pages} further page(s) not OCR'd; "
                     f"raise OCR_MAX_PAGES to include them]")
    return "\n\n".join(pages)


def file_to_text(path: Path, ocr: bool = False) -> dict:
    """Extract text from a downloaded file.

    Always returns a dict with a 'status' so the JSON records *why* a file has
    no text -- unsupported format, missing library, or a scan with no text
    layer. Silence is indistinguishable from an empty document otherwise.
    """
    suffix = path.suffix.lower()
    if suffix not in _EXTRACTORS:
        return {"status": "unsupported", "extractor": None, "text": None}

    library, extractor = _EXTRACTORS[suffix]
    try:
        text = (extractor(path) or "").strip()
    except ImportError:
        return {
            "status": "missing_library",
            "extractor": library,
            "text": None,
            "note": f"pip install {library}",
        }
    except Exception as exc:  # corrupt file, encrypted PDF, etc.
        return {"status": "error", "extractor": library, "text": None, "note": str(exc)[:200]}

    if not text and suffix == ".pdf" and ocr:
        try:
            text = (_ocr_pdf(path) or "").strip()
        except ImportError:
            return {
                "status": "empty", "extractor": library, "text": None,
                "note": "no text layer; OCR needs macOS + "
                        "pip install pyobjc-framework-Vision pyobjc-framework-Quartz",
            }
        except Exception as exc:
            return {"status": "error", "extractor": "ocr", "text": None,
                    "note": f"OCR failed: {str(exc)[:160]}"}
        if text:
            return {
                "status": "ok", "extractor": "macos-vision-ocr", "chars": len(text),
                "truncated": len(text) > MAX_EXTRACT_CHARS,
                "text": text[:MAX_EXTRACT_CHARS],
                "note": "recovered by OCR; no text layer in the original",
            }

    if not text:
        note = ("no text layer (likely a scan; re-run with --ocr)" if suffix == ".pdf"
                else "empty")
        return {"status": "empty", "extractor": library, "text": None, "note": note}

    truncated = len(text) > MAX_EXTRACT_CHARS
    return {
        "status": "ok",
        "extractor": library,
        "chars": len(text),
        "truncated": truncated,
        "text": text[:MAX_EXTRACT_CHARS],
    }
